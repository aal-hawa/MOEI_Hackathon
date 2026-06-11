"""PowerPoint export (python-pptx) — an EXECUTIVE intelligence deck.

KPI cards, native multi-year LINE chart, trade BAR charts, composition DONUT,
plus the full intelligence layer: executive read, talking points, opportunities
vs risks, predictive outlook, the Council verdict and the UAE angle — all built
ONLY from the verified, stored dossier (numbers never come from a model).
All text is cleaned of markdown/citation artifacts before rendering.
"""
from __future__ import annotations

import re
from io import BytesIO
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..timeutils import now_uae
from ..dashboard import parse_shares
from .pdf import FIELD_AR

GOLD = RGBColor(0x9C, 0x7A, 0x2D)
GOLD_L = RGBColor(0xB8, 0x9A, 0x52)
GOLD_PALE = RGBColor(0xF3, 0xEC, 0xDC)
INK = RGBColor(0x2B, 0x2B, 0x2B)
CREAM = RGBColor(0xF7, 0xF2, 0xE8)
MUTED = RGBColor(0x6F, 0x6A, 0x60)
TEAL = RGBColor(0x2C, 0x7A, 0x6B)
TEAL_PALE = RGBColor(0xEA, 0xF4, 0xF1)
BRICK = RGBColor(0xA6, 0x49, 0x2F)
BRICK_PALE = RGBColor(0xF8, 0xEC, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALETTE = [GOLD, TEAL, GOLD_L, BRICK, RGBColor(0xC8, 0xA8, 0x4B),
           RGBColor(0x7A, 0x8C, 0x5A), RGBColor(0x3D, 0x6E, 0x8C)]

UI = {
    "en": {"brief": "Executive Intelligence Briefing", "ministry": "Ministry of Energy & Infrastructure",
           "generated": "Generated", "verified": "every figure is source-verified · nothing invented",
           "kpis": "Key indicators", "trends": "Multi-year trajectory",
           "partners": "Top trade partners (% share)", "exp": "Exports", "imp": "Imports",
           "composition": "Composition", "read": "The executive read",
           "talking": "Talking points", "opps": "Opportunities", "risks": "Risks",
           "oppsrisks": "Opportunities & risks (analysis)", "predictive": "Predictive outlook (projection)",
           "council": "The Council — verdict", "dissent": "Dissent", "uae": "The UAE angle",
           "coverage": "data coverage"},
    "ar": {"brief": "الإحاطة الاستخباراتية التنفيذية", "ministry": "وزارة الطاقة والبنية التحتية",
           "generated": "أُنشئ في", "verified": "كل رقم مُتحقَّق من مصدره · لا اختلاق",
           "kpis": "المؤشرات الرئيسية", "trends": "المسار متعدد السنوات",
           "partners": "أهم الشركاء التجاريين (٪)", "exp": "صادرات", "imp": "واردات",
           "composition": "التركيبة", "read": "الخلاصة التنفيذية",
           "talking": "نقاط الحديث", "opps": "الفرص", "risks": "المخاطر",
           "oppsrisks": "الفرص والمخاطر (تحليل)", "predictive": "التوقعات المستقبلية (إسقاط)",
           "council": "المجلس — الحكم", "dissent": "رأي مخالف", "uae": "زاوية الإمارات",
           "coverage": "تغطية البيانات"},
}
KPIS = {
    "en": [("gdp_nominal_usd", "GDP (nominal)"), ("gdp_per_capita_usd", "GDP / capita"),
           ("gdp_growth_pct", "GDP growth"), ("unemployment_pct", "Unemployment"),
           ("inflation_pct", "Inflation"), ("fdi_inflow_usd", "FDI inflow")],
    "ar": [("gdp_nominal_usd", "الناتج المحلي"), ("gdp_per_capita_usd", "نصيب الفرد"),
           ("gdp_growth_pct", "نمو الناتج"), ("unemployment_pct", "البطالة"),
           ("inflation_pct", "التضخّم"), ("fdi_inflow_usd", "الاستثمار الأجنبي")],
}
UAE_FIELDS = {
    "en": [("uae_investments", "UAE investments"), ("uae_bilateral_trade", "Bilateral trade"),
           ("uae_bilateral_agreements", "Agreements / MoUs"), ("uae_ambassadors", "Ambassadors"),
           ("uae_cooperation_areas", "Cooperation areas"), ("uae_recent_visits", "Recent visits")],
    "ar": [("uae_investments", "الاستثمارات الإماراتية"), ("uae_bilateral_trade", "التجارة الثنائية"),
           ("uae_bilateral_agreements", "الاتفاقيات"), ("uae_ambassadors", "السفراء"),
           ("uae_cooperation_areas", "مجالات التعاون"), ("uae_recent_visits", "الزيارات الأخيرة")],
}


# ── text cleaning (markdown / citation artifacts never reach a slide) ──────────
def _clean(s: Optional[str]) -> str:
    if not s:
        return ""
    s = str(s)
    s = re.sub(r"\(\s*\[[^\]]*\]\([^)]*\)\s*\)", "", s)          # ([label](url))
    s = re.sub(r"\[([^\]]*)\]\((?:https?:)[^)]*\)", r"\1", s)    # [label](url) -> label
    s = re.sub(r"\(\s*https?://[^)]*\)", "", s)                   # (https://…)
    s = s.replace("**", "").replace("__", "")
    s = re.sub(r"^#{1,4}\s*", "", s, flags=re.M)                  # headings
    s = re.sub(r"\?utm_source=\w+", "", s)
    s = re.sub(r"[ \t]{2,}", " ", s)
    return s.strip()


def _lines(s: Optional[str], maxn: Optional[int] = None) -> List[str]:
    out = []
    for ln in _clean(s).splitlines():
        ln = ln.strip().lstrip("•-–—*0123456789.").strip()
        if ln:
            out.append(ln)
    return out[:maxn] if maxn else out


def _split_opp_risk(analysis: Optional[str]) -> Tuple[List[str], List[str]]:
    """Heuristic split of the analyst draft into opportunity vs risk bullets."""
    opp: List[str] = []
    rsk: List[str] = []
    cur = opp
    for ln in _lines(analysis):
        low = ln.lower()
        if len(ln) < 70 and ("risk" in low or "مخاطر" in ln or "تحديات" in ln) and ln.rstrip(":").count(" ") <= 6:
            cur = rsk
            continue
        if len(ln) < 70 and ("opportunit" in low or "فرص" in ln) and ln.rstrip(":").count(" ") <= 6:
            cur = opp
            continue
        if "risk" in low.split(":")[0][:30] and ":" in ln:
            rsk.append(ln)
            continue
        cur.append(ln)
    return opp, rsk


def _parse_council(text: Optional[str]):
    """Parse the council format into (seats, verdict, confidence, dissent)."""
    seats: List[Tuple[str, str]] = []
    verdict = confidence = dissent = ""
    names = ["ECONOMIST", "ENERGY & INFRA STRATEGIST", "GEOPOLITICAL ADVISOR", "RISK OFFICER"]
    cur = None
    for ln in _lines(text):
        upper = ln.upper()
        seat = next((n for n in names if upper.startswith(n)), None)
        if seat:
            seats.append((seat.title(), ln.split(":", 1)[-1].strip()))
            cur = "seat"
        elif upper.startswith("VERDICT"):
            verdict = ln.split(":", 1)[-1].strip(); cur = "verdict"
        elif upper.startswith("CONFIDENCE"):
            confidence = ln.split(":", 1)[-1].strip(); cur = None
        elif upper.startswith("DISSENT"):
            dissent = ln.split(":", 1)[-1].strip(); cur = None
        elif cur == "seat" and seats:
            seats[-1] = (seats[-1][0], (seats[-1][1] + " " + ln).strip())
        elif cur == "verdict":
            verdict = (verdict + " " + ln).strip()
    return seats, verdict, confidence, dissent


# ── drawing helpers ─────────────────────────────────────────────────────────
def _fmt(value, unit):
    if value is None:
        return "—"
    try:
        n = float(value)
        if unit == "USD":
            if abs(n) >= 1e12: return f"${n/1e12:.2f}T"
            if abs(n) >= 1e9: return f"${n/1e9:.1f}B"
            if abs(n) >= 1e6: return f"${n/1e6:.1f}M"
            return f"${n:,.0f}"
        if unit == "%": return f"{n:.1f}%"
        if unit == "people": return f"{n:,.0f}"
    except (ValueError, TypeError):
        pass
    return f"{value}{(' ' + unit) if unit else ''}"


def _txt(slide, left, top, width, height, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text or ""
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb


def _bullets(slide, x, y, w, h, items, size=14, color=INK, align=PP_ALIGN.LEFT,
             bullet="•", space=7):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run(); r.text = f"{bullet}  {it}"
        r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def _panel(slide, x, y, w, h, fill=WHITE, edge=GOLD_L, accent: Optional[RGBColor] = GOLD):
    box = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = edge; box.line.width = Pt(1)
    if accent is not None:
        bar = slide.shapes.add_shape(1, x, y, Inches(0.07), h)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    return box


class _Deck:
    """Slide factory with the brand frame + footer page numbers."""
    def __init__(self, prs, country, ui, ar):
        self.prs = prs; self.country = country; self.ui = ui; self.ar = ar; self.n = 0

    def slide(self, title: Optional[str] = None, subtitle: Optional[str] = None):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = s.background.fill; bg.solid(); bg.fore_color.rgb = CREAM
        bar = s.shapes.add_shape(1, 0, 0, self.prs.slide_width, Inches(0.16))
        bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()
        self.n += 1
        al = PP_ALIGN.RIGHT if self.ar else PP_ALIGN.LEFT
        if title:
            _txt(s, Inches(0.8), Inches(0.42), Inches(11.7), Inches(0.7), title, 27, INK, True, al)
            ln = s.shapes.add_shape(1, Inches(0.8), Inches(1.12), Inches(2.2), Pt(3))
            ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
        if subtitle:
            _txt(s, Inches(0.8), Inches(1.16), Inches(11.7), Inches(0.4), subtitle, 12, MUTED, False, al)
        if self.n > 1:  # footer (skip title slide)
            _txt(s, Inches(0.8), Inches(7.06), Inches(9.0), Inches(0.35),
                 f"{self.country}  ·  {self.ui['ministry']}", 9, MUTED, False,
                 PP_ALIGN.RIGHT if self.ar else PP_ALIGN.LEFT)
            _txt(s, Inches(12.2), Inches(7.06), Inches(0.8), Inches(0.35),
                 str(self.n), 10, GOLD, True, PP_ALIGN.RIGHT)
        return s


def _style_chart(chart, legend=False, palette=True):
    try:
        chart.has_title = False
        chart.has_legend = legend
        if legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
    except Exception:
        pass
    try:
        plot = chart.plots[0]
        plot.has_data_labels = False
        series = plot.series[0]
        if palette:
            for i, pt in enumerate(series.points):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = PALETTE[i % len(PALETTE)]
        else:
            series.format.line.color.rgb = GOLD
            series.format.line.width = Pt(2.25)
    except Exception:
        pass


# ── the deck ─────────────────────────────────────────────────────────────────
def build_pptx(country: str, iso3: str, rows: List[Dict], trends: Optional[Dict] = None,
               trade: Optional[Dict] = None, summary: Optional[str] = None,
               lang: str = "en", analysis: Optional[str] = None,
               talking_points: Optional[str] = None, predictive: Optional[str] = None,
               council: Optional[str] = None, coverage: Optional[float] = None) -> bytes:
    ar = (lang == "ar")
    ui = UI["ar" if ar else "en"]
    al = PP_ALIGN.RIGHT if ar else PP_ALIGN.LEFT
    trends = trends or {}
    trade = trade or {}
    by = {r["field_name"]: r for r in rows}

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    deck = _Deck(prs, country, ui, ar)

    # ── 1) Title ──
    s = deck.slide()
    hexa = s.shapes.add_shape(9, Inches(0.8), Inches(1.7), Inches(0.85), Inches(0.85))  # hexagon
    hexa.fill.solid(); hexa.fill.fore_color.rgb = GOLD; hexa.line.fill.background()
    _txt(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.8), ui["brief"], 20, MUTED, True, al)
    _txt(s, Inches(0.8), Inches(3.3), Inches(11.7), Inches(1.4), country, 56, INK, True, al)
    _txt(s, Inches(0.8), Inches(4.75), Inches(11.7), Inches(0.6), f"{iso3}  ·  {ui['ministry']}", 16, GOLD, True, al)
    cov = f"  ·  {round(coverage*100)}% {ui['coverage']}" if coverage else ""
    _txt(s, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.5),
         f"{ui['generated']} {now_uae():%d %b %Y}  ·  {ui['verified']}{cov}", 11, MUTED, False, al)

    # ── 2) Executive read ──
    read_lines = _lines(summary)
    if read_lines:
        s = deck.slide(ui["read"])
        body = " ".join(read_lines)
        paras = [p.strip() for p in re.split(r"(?<=[.؟!?])\s+(?=[A-Z«؀-ۿ])", body) if p.strip()]
        # group sentences into ~4 readable paragraphs
        chunks: List[str] = []
        cur = ""
        for p in paras:
            if len(cur) + len(p) < 420:
                cur = (cur + " " + p).strip()
            else:
                chunks.append(cur); cur = p
        if cur:
            chunks.append(cur)
        _panel(s, Inches(0.8), Inches(1.45), Inches(11.73), Inches(5.4), fill=WHITE)
        _bullets(s, Inches(1.05), Inches(1.7), Inches(11.2), Inches(4.9),
                 chunks[:5], size=15, align=al, bullet="—", space=12)

    # ── 3) KPI cards ──
    cards = [(k, lbl) for k, lbl in KPIS["ar" if ar else "en"]
             if by.get(k) and by[k].get("value") is not None]
    if cards:
        s = deck.slide(ui["kpis"])
        x0, y0, w, h = Inches(0.8), Inches(1.6), Inches(3.85), Inches(2.2)
        gapx, gapy = Inches(0.18), Inches(0.3)
        for i, (k, lbl) in enumerate(cards[:6]):
            r = by[k]
            col, row = i % 3, i // 3
            x = x0 + col * (w + gapx); y = y0 + row * (h + gapy)
            card = _panel(s, x, y, w, h)
            tf = card.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = al
            run = p.add_run(); run.text = _fmt(r["value"], r["unit"])
            run.font.size = Pt(30); run.font.bold = True; run.font.color.rgb = INK
            p2 = tf.add_paragraph(); p2.alignment = al
            r2 = p2.add_run(); r2.text = lbl; r2.font.size = Pt(13); r2.font.color.rgb = MUTED
            tr = trends.get(k)
            if tr:
                p3 = tf.add_paragraph(); p3.alignment = al
                move = (f"{tr.get('delta_pp', 0):+} pp" if tr.get("unit") == "%"
                        else f"{tr.get('cagr_pct', 0)}%/yr")
                good = (tr.get("direction") == "up") == bool(tr.get("good_up")) or tr.get("direction") == "flat"
                r3 = p3.add_run(); r3.text = f"{'▲' if tr.get('direction')=='up' else '▼'} {move} / {tr.get('span_years','')}y"
                r3.font.size = Pt(12); r3.font.bold = True
                r3.font.color.rgb = TEAL if good else BRICK
            src = r.get("source_name") or ""
            if src:
                p4 = tf.add_paragraph(); p4.alignment = al
                r4 = p4.add_run(); r4.text = f"{src}{(' · ' + str(r.get('as_of_date'))) if r.get('as_of_date') else ''}"
                r4.font.size = Pt(8.5); r4.font.color.rgb = MUTED

    # ── 4) Multi-year trajectory (native LINE) ──
    line_keys = [k for k in ("gdp_per_capita_usd", "gdp_growth_pct", "unemployment_pct")
                 if trends.get(k) and trends[k].get("spark")]
    if line_keys:
        s = deck.slide(ui["trends"])
        key = line_keys[0]
        tr = trends[key]
        vals = tr["spark"]; yr = tr.get("latest_year") or 0
        cats = [str(yr - (len(vals) - 1 - i)) for i in range(len(vals))]
        cd = CategoryChartData(); cd.categories = cats
        lbl = (FIELD_AR.get(key) if ar else None) or key.replace("_", " ")
        cd.add_series(lbl, vals)
        gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.8), Inches(1.5),
                                Inches(11.7), Inches(5.2), cd)
        _style_chart(gf.chart, legend=False, palette=False)

    # ── 5) Top trade partners (native BAR) ──
    def _partner_chart(title, key):
        arr = (trade.get(key) or [])[:6]
        if not arr:
            return
        s = deck.slide(title)
        cd = CategoryChartData(); cd.categories = [a["name"] for a in arr]
        cd.add_series("%", [a.get("share_pct") or 0 for a in arr])
        gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.5),
                                Inches(11.7), Inches(5.2), cd)
        _style_chart(gf.chart, legend=False, palette=True)

    _partner_chart(f"{ui['partners']} — {ui['exp']}", "export_partners")
    _partner_chart(f"{ui['partners']} — {ui['imp']}", "import_partners")

    # ── 6) Composition donut ──
    comp = None; comp_title = ui["composition"]
    if trade.get("export_goods"):
        comp = [(g["name"], g.get("share_pct") or 0) for g in trade["export_goods"][:6]]
        comp_title = f"{ui['composition']} — {ui['exp']}"
    elif by.get("gdp_by_sector") and parse_shares(by["gdp_by_sector"]["value"]):
        comp = [(x["name"], x["value"]) for x in parse_shares(by["gdp_by_sector"]["value"])]
        comp_title = (FIELD_AR.get("gdp_by_sector") if ar else "GDP by sector")
    if comp:
        s = deck.slide(comp_title)
        cd = CategoryChartData(); cd.categories = [c[0] for c in comp]
        cd.add_series("share", [c[1] for c in comp])
        gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(2.5), Inches(1.4),
                                Inches(8.3), Inches(5.3), cd)
        _style_chart(gf.chart, legend=True, palette=True)

    # ── 7) The UAE angle ──
    uae_facts = [(lbl, _clean(str(by[k]["value"]))) for k, lbl in UAE_FIELDS["ar" if ar else "en"]
                 if by.get(k) and by[k].get("value")]
    if uae_facts:
        s = deck.slide(ui["uae"])
        x0, y0, w, h = Inches(0.8), Inches(1.5), Inches(5.85), Inches(1.72)
        for i, (lbl, val) in enumerate(uae_facts[:6]):
            col, row = i % 2, i // 2
            x = x0 + col * (w + Inches(0.2)); y = y0 + row * (h + Inches(0.18))
            card = _panel(s, x, y, w, h, accent=TEAL)
            tf = card.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = al
            r1 = p.add_run(); r1.text = lbl
            r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = GOLD
            p2 = tf.add_paragraph(); p2.alignment = al
            r2 = p2.add_run(); r2.text = val[:220] + ("…" if len(val) > 220 else "")
            r2.font.size = Pt(11.5); r2.font.color.rgb = INK

    # ── 8) Opportunities & risks ──
    opp, rsk = _split_opp_risk(analysis)
    if opp or rsk:
        s = deck.slide(ui["oppsrisks"])
        half = Inches(5.85)
        _panel(s, Inches(0.8), Inches(1.5), half, Inches(5.3), fill=TEAL_PALE, edge=TEAL, accent=TEAL)
        _txt(s, Inches(1.05), Inches(1.65), Inches(5.3), Inches(0.4), ui["opps"], 15, TEAL, True, al)
        _bullets(s, Inches(1.05), Inches(2.15), Inches(5.4), Inches(4.5), [o[:200] for o in opp[:6]],
                 size=11.5, align=al)
        _panel(s, Inches(6.85), Inches(1.5), half, Inches(5.3), fill=BRICK_PALE, edge=BRICK, accent=BRICK)
        _txt(s, Inches(7.1), Inches(1.65), Inches(5.3), Inches(0.4), ui["risks"], 15, BRICK, True, al)
        _bullets(s, Inches(7.1), Inches(2.15), Inches(5.4), Inches(4.5), [r[:200] for r in (rsk or opp[6:])[:6]],
                 size=11.5, align=al)

    # ── 9) Predictive outlook ──
    pred = _lines(predictive, 7)
    if pred:
        s = deck.slide(ui["predictive"])
        _panel(s, Inches(0.8), Inches(1.45), Inches(11.73), Inches(5.4))
        _bullets(s, Inches(1.05), Inches(1.7), Inches(11.2), Inches(4.9),
                 [p[:260] for p in pred], size=13.5, align=al, space=10)

    # ── 10) Talking points ──
    talk = _lines(talking_points, 9)
    if talk:
        s = deck.slide(ui["talking"])
        _panel(s, Inches(0.8), Inches(1.45), Inches(11.73), Inches(5.4))
        _bullets(s, Inches(1.05), Inches(1.7), Inches(11.2), Inches(4.9),
                 [t[:240] for t in talk], size=13, align=al, bullet="🗣", space=9)

    # ── 11) Council verdict ──
    seats, verdict, confidence, dissent = _parse_council(council)
    if verdict or seats:
        s = deck.slide(ui["council"])
        if seats:
            x0, y0, w, h = Inches(0.8), Inches(1.45), Inches(2.86), Inches(2.1)
            for i, (name, body) in enumerate(seats[:4]):
                x = x0 + i * (w + Inches(0.12))
                card = _panel(s, x, y0, w, h, accent=None)
                tf = card.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = al
                r1 = p.add_run(); r1.text = name
                r1.font.size = Pt(10.5); r1.font.bold = True; r1.font.color.rgb = GOLD
                p2 = tf.add_paragraph(); p2.alignment = al
                r2 = p2.add_run(); r2.text = body[:230] + ("…" if len(body) > 230 else "")
                r2.font.size = Pt(9.5); r2.font.color.rgb = INK
        if verdict:
            vp = _panel(s, Inches(0.8), Inches(3.75), Inches(11.73), Inches(2.2),
                        fill=GOLD_PALE, edge=GOLD, accent=GOLD)
            tf = vp.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = al
            r1 = p.add_run(); r1.text = verdict[:520]
            r1.font.size = Pt(13.5); r1.font.bold = True; r1.font.color.rgb = INK
        foot = []
        if confidence:
            foot.append(("🎯 " + confidence[:120], TEAL))
        if dissent and dissent.lower() not in ("none", "لا يوجد"):
            foot.append((f"{ui['dissent']}: {dissent[:160]}", BRICK))
        yy = Inches(6.1)
        for txt_, col_ in foot:
            _txt(s, Inches(0.8), yy, Inches(11.7), Inches(0.4), txt_, 11.5, col_, True, al)
            yy += Inches(0.42)

    bio = BytesIO(); prs.save(bio)
    return bio.getvalue()
